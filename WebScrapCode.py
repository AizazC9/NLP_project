import requests
import re
import pandas as pd
import nltk
from bs4 import BeautifulSoup
from nltk.stem import PorterStemmer
from nltk.stem import WordNetLemmatizer
from textstat import flesch_kincaid_grade
import os
from PIL import Image
import pytesseract
from io import BytesIO
import PyPDF2
import io
from pptx import Presentation
from urllib.parse import urljoin
from PyPDF2 import PdfReader

# Function to extract text from image
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
def text_from_image(image_url):
    if not image_url.lower().endswith(('.png', '.jpg', '.jpeg')):  # Skip non-standard formats
        print(f"Skipping non-standard image format: {image_url}")
        return ""
    
    try:
        response = requests.get(image_url)
        img = Image.open(BytesIO(response.content))
        text = pytesseract.image_to_string(img)
        return text
    except Exception as e:
        print(f"Error processing image URL {image_url}: {e}")
        return ""

# Function to scrape tables
def scrape_table(soup):
    tables = soup.find_all('table')
    extracted_tables = []
    for table in tables:
        headers = [header.text.strip() for header in table.find_all('th')]
        if not headers:  # If no headers, try to infer them from the first row
            first_row = table.find('tr')
            if first_row:
                headers = [cell.text.strip() for cell in first_row.find_all('td')]
        
        rows = []
        for row in table.find_all('tr'):
            cells = row.find_all(['td', 'th'])
            row_data = [cell.text.strip() for cell in cells]
            if len(row_data) == len(headers):
                rows.append(row_data)
            else:
                # Skip rows that don't match header count to avoid DataFrame creation error
                continue
        
        if headers and rows:
            df = pd.DataFrame(rows, columns=headers)
            extracted_tables.append(df)
    return extracted_tables

# Function to extract text from PDF
def text_from_pdf(pdf_url):
    try:
        response = requests.get(pdf_url)
        pdf_file = io.BytesIO(response.content)

        reader = PdfReader(pdf_file)
        num_pages = len(reader.pages)
        pdf_text = ''

        for page in range(num_pages):
            pdf_text += reader.pages[page].extract_text()

        return pdf_text
    except Exception as e:
        print(f"Error processing PDF URL {pdf_url}: {e}")
        return ""

# Function to extract text from PPTX
def text_from_pptx(pptx_url):
    try:
        response = requests.get(pptx_url)
        pptx_file = io.BytesIO(response.content)
        prs = Presentation(pptx_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error processing PPTX URL {pptx_url}: {e}")
        return ""

# Function to process different file types (PDF and PPTX)
def process_file(file_url, file_text):
    if file_url.endswith('.pptx'):
        file_text += text_from_pptx(file_url) + "\n"
    elif file_url.endswith('.pdf'):
        file_text += text_from_pdf(file_url) + "\n"
    return file_text

# Function to detect PDF and PPTX URLs
def find_pdf_pptx_urls(soup, base_url):
    file_urls = []
    for link in soup.find_all('a', href=True):
        href = link['href']
        if href.endswith('.pdf') or href.endswith('.pptx'):
            absolute_url = urljoin(base_url, href)
            file_urls.append(absolute_url)
    return file_urls

# Function to scrape code snippets
def scrape_code(soup):
    code_blocks = soup.find_all('pre')
    extracted_code = []
    for block in code_blocks:
        code = ''.join([code.text for code in block.find_all('code')])
        extracted_code.append(code)
    return extracted_code

# Function to append or concatenate data in CSV
def update_or_append_csv(filename, new_row):
    if os.path.exists(filename):
        df = pd.read_csv(filename)
        topic_index = df[df['topic'] == new_row['topic'][0]].index
        if len(topic_index) > 0:
            index = topic_index[0]
            for column in new_row.columns:
                if column != 'topic':
                    if column in ['vocabulary size', 'word count']:
                        df.at[index, column] = df.at[index, column] + new_row.at[0, column]
                    elif column == 'type':
                        if df.at[index, column] != new_row.at[0, column]:
                            df.at[index, column] = new_row.at[0, column]
                    else:
                        df.at[index, column] = str(df.at[index, column]) + ', ' + str(new_row.at[0, column])
        else:
            df = pd.concat([df, new_row], ignore_index=True)
    else:
        df = new_row
    df.to_csv(filename, index=False)

def update_readability_csv(filename, new_row):
    if os.path.exists(filename):
        df = pd.read_csv(filename)
        topic_index = df[df['topic'] == new_row['topic'][0]].index
        if len(topic_index) > 0:
            index = topic_index[0]
            if df.at[index, 'readability'] != new_row.at[0, 'readability']:
                df.at[index, 'readability'] = new_row.at[0, 'readability']
        else:
            df = pd.concat([df, new_row], ignore_index=True)
    else:
        df = new_row
    df.to_csv(filename, index=False)

#  75 subtopics of Web Development
topics = [
    "HTML Basics","CSS for Styling","JavaScript Fundamentals","Responsive Web Design","React","Back-End Development","Server-Side Scripting (e.g., Node.js, Python, PHP)","Database Management (e.g., SQL, MongoDB)","RESTful APIs","JSON and XML",
    "Web Hosting and Deployment","Version Control Systems (e.g., Git)","Content Management Systems (e.g., WordPress, Drupal)","Search Engine Optimization (SEO)","Web Security Practices","Cross-Browser Compatibility","UI/UX Design Principles","Graphic Design for Web",
    "CSS Preprocessors (e.g., SASS, LESS)","JavaScript Libraries (e.g., jQuery)","AJAX for Asynchronous Requests","Web Accessibility Standards",
    "Performance Optimization","Progressive Web Apps (PWAs)","E-commerce Development","Web Analytics","Mobile-First Design","Web Frameworks (e.g., Django, Flask, Express)",
    "Full Stack Development","Testing and Debugging (e.g., unit tests, integration tests)","Websockets for Real-Time Communication","Cloud Computing Services (e.g., AWS, Azure)","DevOps Practices for Web Development",
    "Continuous Integration/Continuous Deployment (CI/CD)","Web Development Tools and IDEs","Agile and Scrum Methodologies","Microservices Architecture",
    "API Design and Development","Containerization (e.g., Docker)","Scalability in Web Applications","Content Delivery Networks (CDNs)","SSL Certificates and HTTPS","Domain Name System (DNS) Management",
    "Load Balancing Techniques","Web Caching Strategies","Payment Gateway Integration","Email Services in Web Development","Social Media Integration",
    "Web Scraping Techniques","Data Visualization in Web","Animation and Motion in Web Design",
    "Typography in Web Design","Color Theory for Web","User Authentication and Authorization","File Upload and Management","Custom Web Components",
    "GraphQL","Block Chain","Internet Protocols and Standards",
    "MVC Architecture in Web Development","Web Application Security Threats","Cross-Origin Resource Sharing (CORS)","Internationalization and Localization",
    "Automated Web Testing Tools","Virtual DOM in Web Development","State Management in Web Apps","Serverless Architecture","Web APIs (Geolocation, Web Storage, etc.)","Browser Developer Tools",
    "Code Quality and Refactoring","Progressive Enhancement and Graceful Degradation","Web Development Project Management","User Feedback and Usability Testing",
    "Accessibility Testing","Ethical and Legal Aspects of Web Development"
]

#  75 topics types
types = [
    "Front-End Development",  # HTML Basics
    "Front-End Development",  # CSS for Styling
    "Front-End Development",  # JavaScript Fundamentals
    "Front-End Development",  # Responsive Web Design
    "Front-End Development",  # Front-End Frameworks (e.g., React, Angular, Vue)
    "Back-End Development",  # Back-End Development
    "Back-End Development",  # Server-Side Scripting (e.g., Node.js, Python, PHP)
    "Back-End Development",  # Database Management (e.g., SQL, MongoDB)
    "Back-End Development",  # RESTful APIs
    "Data Handling",  # JSON and XML
    "Deployment",  # Web Hosting and Deployment
    "Development Tools",  # Version Control Systems (e.g., Git)
    "Back-End Development",  # Content Management Systems (e.g., WordPress, Drupal)
    "SEO and Marketing",  # Search Engine Optimization (SEO)
    "Security",  # Web Security Practices
    "Front-End Development",  # Cross-Browser Compatibility
    "UI/UX Design",  # UI/UX Design Principles
    "UI/UX Design",  # Graphic Design for Web
    "Front-End Development",  # CSS Preprocessors (e.g., SASS, LESS)
    "Front-End Development",  # JavaScript Libraries (e.g., jQuery)
    "Front-End Development",  # AJAX for Asynchronous Requests
    "UI/UX Design",  # Web Accessibility Standards
    "Optimization",  # Performance Optimization
    "Front-End Development",  # Progressive Web Apps (PWAs)
    "Full Stack Development",  # E-commerce Development
    "Analytics",  # Web Analytics
    "UI/UX Design",  # Mobile-First Design
    "Back-End Development",  # Web Frameworks (e.g., Django, Flask, Express)
    "Full Stack Development",  # Full Stack Development
    "Testing and Debugging",  # Testing and Debugging (e.g., unit tests, integration tests)
    "Front-End Development",  # Websockets for Real-Time Communication
    "Cloud Computing",  # Cloud Computing Services (e.g., AWS, Azure)
    "DevOps",  # DevOps Practices for Web Development
    "DevOps",  # Continuous Integration/Continuous Deployment (CI/CD)
    "Development Tools",  # Web Development Tools and IDEs
    "Project Management",  # Agile and Scrum Methodologies
    "Back-End Development",  # Microservices Architecture
    "Back-End Development",  # API Design and Development
    "DevOps",  # Containerization (e.g., Docker)
    "Back-End Development",  # Scalability in Web Applications
    "Infrastructure",  # Content Delivery Networks (CDNs)
    "Security",  # SSL Certificates and HTTPS
    "Infrastructure",  # Domain Name System (DNS) Management
    "Infrastructure",  # Load Balancing Techniques
    "Infrastructure",  # Web Caching Strategies
    "E-commerce",  # Payment Gateway Integration
    "Marketing",  # Email Services in Web Development
    "Social Media",  # Social Media Integration
    "Data Handling",  # Web Scraping Techniques
    "Data Visualization",  # Data Visualization in Web
    "UI/UX Design",  # Animation and Motion in Web Design
    "UI/UX Design",  # Typography in Web Design
    "UI/UX Design",  # Color Theory for Web
    "Security",  # User Authentication and Authorization
    "Back-End Development",  # File Upload and Management
    "Front-End Development",  # Custom Web Components
    "Back-End Development",  # GraphQL
    "Infrastructure",  # BlockChain
    "Networking",  # Internet Protocols and Standards
    "Back-End Development",  # MVC Architecture in Web Development
    "Security",  # Web Application Security Threats
    "Networking",  # Cross-Origin Resource Sharing (CORS)
    "Localization",  # Internationalization and Localization
    "Testing and Debugging",  # Automated Web Testing Tools
    "Front-End Development",  # Virtual DOM in Web Development
    "Front-End Development",  # State Management in Web Apps
    "Cloud Computing",  # Serverless Architecture
    "APIs",  # Web APIs (Geolocation, Web Storage, etc.)
    "Development Tools",  # Browser Developer Tools
    "Code Quality",  # Code Quality and Refactoring
    "UI/UX Design",  # Progressive Enhancement and Graceful Degradation
    "Project Management",  # Web Development Project Management
    "UI/UX Design",  # User Feedback and Usability Testing
    "UI/UX Design",  # Accessibility Testing
    "Ethics",  # Ethical and Legal Aspects of Web Development
]

# URL to scrape
url = "https://medium.com/@corewave/ethical-considerations-in-web-development-7015bac034cb"

# HTTP GET request to the URL
response = requests.get(url)

# topic for data
index = 74 
topic = topics[index]
typ = types[index]

# Checking if the request was successful (status code 200)
if response.status_code == 200:
    # Parsing the HTML content of the page
    soup = BeautifulSoup(response.content, "html.parser")

    # Scraping text from images
    image_urls = [img['src'] for img in soup.find_all('img') if img.get('src')]
    image_texts = [text_from_image(url) for url in image_urls]

    # Scraping data from tables
    tables = scrape_table(soup)

    # Scraping code snippets
    code_snippets = scrape_code(soup)

    # Finding and scraping PDFs and PPTXs
    file_urls = find_pdf_pptx_urls(soup, url)
    file_text = ""
    for file_url in file_urls:
        file_text = process_file(file_url, file_text)

    # Finding the main content element and extracting text
    main_content = soup.find("article", {"id": ""})
   
    data = ''
    if main_content:
        data = main_content.get_text()
    
        # Add text from images, tables, and code snippets to the main data
        data = data + "\n".join(image_texts) + "\n".join([table.to_string() for table in tables]) + "\n".join(code_snippets) + file_text

        # Removing extra whitespaces and special characters using regular expressions
        data = re.sub(r'\s+', ' ', data)
        data = re.sub(r'[^\w\s]', '', data)

        # Calculating vocabulary size and word count
        words = data.split()
        vocabulary_size = len(set(words))
        word_count = len(words)

        # Preparing data for corpus.csv
        corpus_df = pd.DataFrame({'topic': [topic], 'content': [data], 'URL': [url], 'type': [typ], 
                                    'vocabulary size': [vocabulary_size], 'word count': [word_count]})

        # Calculating Flesch-Kincaid grade level
        flesch_grade_level = flesch_kincaid_grade(data)

        # Determining readability
        readability = 'easy' if flesch_grade_level <= 8 else 'difficult'

        # Preparing data for readability.csv
        readability_df = pd.DataFrame({'topic': [topic], 'readability': [readability]})

        # Tokenizing the data
        tokens = nltk.word_tokenize(data)

        # Stemming and Lemmatization
        stemmer = PorterStemmer()
        lemmatizer = WordNetLemmatizer()
        stems = [stemmer.stem(word) for word in tokens]
        lemmas = [lemmatizer.lemmatize(word) for word in tokens]

        # Preparing data for tokens.csv and other CSVs
        tokens_df = pd.DataFrame({'topic': [topic], 'stems': [' '.join(stems)], 'lemmas': [' '.join(lemmas)]})
        unique_words = set(words)
        vocabulary_df = pd.DataFrame({'topic': [topic], 'unique words': [' '.join(unique_words)]})
        pos_tags = nltk.pos_tag(tokens)
        pos_df = pd.DataFrame({'topic': [topic], 'data with pos tags': [pos_tags]})

        # Updating or appending data in CSV files
        update_or_append_csv('corpus.csv', corpus_df)
        update_or_append_csv('tokens.csv', tokens_df)
        update_or_append_csv('vocabulary.csv', vocabulary_df)
        update_or_append_csv('pos.csv', pos_df)

        # Updating readability.csv
        update_readability_csv('readability.csv', readability_df)

        print("Data has been successfully scraped, cleaned, and saved.")
    else:
        print("Main content element not found on the page.")
else:
    print("Failed to retrieve the page.")
